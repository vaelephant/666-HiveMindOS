import hashlib
import uuid
from pathlib import Path

from knowledge_base import config
from server.logging_config import get_logger
from knowledge_base.core.compiler.entity_extractor import extract_entities
from knowledge_base.core.compiler.workflow_extractor import extract_workflows
from knowledge_base.core.compiler.entity_resolver import EntityResolver
from knowledge_base.core.compiler.wiki_merger import (
    upsert_entity_page, upsert_workflow_page, upsert_rule_page,
)
from knowledge_base.core.wiki.wiki_manager import WikiManager
from knowledge_base.core.wiki import wiki_meta
from knowledge_base.core.graph.memory_graph import MemoryGraph
from knowledge_base.core.services.candidate_service import enqueue_from_ingest_compile
from knowledge_base.models.entity import Entity, Relation

log = get_logger("hivemind.agent.ingest")


class IngestAgent:
    def __init__(self, wiki: WikiManager, graph: MemoryGraph):
        self.wiki = wiki
        self.graph = graph

    def run(
        self,
        file_path: Path,
        org_id: str,
        source_type: str = "pdf",
        source_id: str | None = None,
    ) -> dict:
        source_filename = file_path.name

        # ── 1. 文本提取 ───────────────────────────────────────────────────────
        log.info("extracting text  file=%s  type=%s", source_filename, source_type)
        content = self._extract_text(file_path, source_type)
        content_hash = hashlib.sha256(content.encode()).hexdigest()
        log.info("text extracted   chars=%d  hash=%s…", len(content), content_hash[:8])

        # ── 2. LLM 提取结构化数据 ─────────────────────────────────────────────
        log.info("extracting entities via LLM…")
        raw_entities = extract_entities(content)
        log.info("entities found   count=%d", len(raw_entities))

        log.info("extracting workflows via LLM…")
        workflows_data = extract_workflows(content)
        log.info("workflows found  count=%d  rules=%d",
                 len(workflows_data.get("workflows", [])),
                 len(workflows_data.get("rules", [])))

        # ── 3. 实体消歧 / 合并 ────────────────────────────────────────────────
        resolver = EntityResolver(self.graph)
        resolved = resolver.resolve_all(org_id, raw_entities, source_filename)

        new_count = sum(1 for r in resolved if r.is_new)
        merged_count = sum(1 for r in resolved if not r.is_new)
        conflict_total = sum(len(r.conflicts) for r in resolved)
        log.info("resolved  new=%d  merged=%d  conflicts=%d", new_count, merged_count, conflict_total)

        # ── 4. Wiki 增量更新 + 图谱写入 ───────────────────────────────────────
        pages: list[str] = []
        mirror_items: list[dict] = []
        name_to_id: dict[str, str] = {}

        for raw, r in zip(raw_entities, resolved):
            wiki_path = upsert_entity_page(self.wiki.root, org_id, r, source_filename)
            pages.append(wiki_path)
            mirror_items.append({
                "wiki_path": wiki_path,
                "category": "entity",
                "title": r.name,
                "content": r.description or "",
            })
            name_to_id[r.name] = r.entity_id

            conflict_rows = [
                {
                    "field": c.field,
                    "existing_value": c.existing_value,
                    "new_value": c.new_value,
                    "source": c.new_source,
                }
                for c in r.conflicts
            ]
            wiki_meta.record_entity_compile(
                self.wiki.root,
                org_id,
                wiki_path,
                source_id=source_id,
                source_filename=source_filename,
                source_type=source_type,
                entity_type=r.entity_type,
                description=r.description,
                attributes=r.all_attributes,
                attribute_provenance=raw.get("attribute_provenance", {}),
                relations=r.relations,
                conflicts=conflict_rows,
                is_new=r.is_new,
                new_attributes=r.new_attributes,
            )

            self.graph.upsert_entity(Entity(
                id=r.entity_id,
                org_id=org_id,
                name=r.name,
                entity_type=r.entity_type,
                wiki_path=wiki_path,
                attributes=r.all_attributes,
            ))

        # ── 5. 关系合并写入 ───────────────────────────────────────────────────
        rel_count = 0
        for r in resolved:
            src_id = name_to_id.get(r.name)
            if not src_id:
                continue
            for rel in r.relations:
                tgt_id = name_to_id.get(rel.get("target", ""))
                if not tgt_id:
                    continue
                self.graph.add_relation(Relation(
                    id=str(uuid.uuid4()),
                    org_id=org_id,
                    source_entity_id=src_id,
                    target_entity_id=tgt_id,
                    relation_type=rel.get("type", "related"),
                    weight=1.0,
                ))
                rel_count += 1
        log.info("relations written  count=%d", rel_count)

        # ── 6. 流程 & 规则写入 ────────────────────────────────────────────────
        for wf in workflows_data.get("workflows", []):
            wf_path = upsert_workflow_page(self.wiki.root, org_id, wf, source_filename)
            pages.append(wf_path)
            mirror_items.append({
                "wiki_path": wf_path,
                "category": "workflow",
                "title": wf["name"],
                "content": "\n".join(wf.get("steps", [])),
            })
            wiki_meta.record_workflow_compile(
                self.wiki.root, org_id, wf_path,
                source_id=source_id,
                source_filename=source_filename,
                source_type=source_type,
                workflow=wf,
            )

        for rule in workflows_data.get("rules", []):
            rule_path = upsert_rule_page(self.wiki.root, org_id, rule, source_filename)
            pages.append(rule_path)
            mirror_items.append({
                "wiki_path": rule_path,
                "category": "rule",
                "title": rule["name"],
                "content": rule.get("condition", "") or rule.get("action", ""),
            })
            wiki_meta.record_rule_compile(
                self.wiki.root, org_id, rule_path,
                source_id=source_id,
                source_filename=source_filename,
                source_type=source_type,
                rule=rule,
            )

        # ── 7. 候选池 mirror（ingest 双写）────────────────────────────────────
        enqueue_from_ingest_compile(org_id, source_id, source_filename, mirror_items)

        # ── 8. 索引更新 ───────────────────────────────────────────────────────
        self.wiki.update_index(org_id)
        log.info("wiki index updated  org=%s  total_pages=%d", org_id, len(pages))

        return {
            "content_hash": content_hash,
            "entities_extracted": len(raw_entities),
            "entities_new": new_count,
            "entities_merged": merged_count,
            "conflicts_detected": conflict_total,
            "workflows_extracted": len(workflows_data.get("workflows", [])),
            "rules_extracted": len(workflows_data.get("rules", [])),
            "wiki_pages_created": len(pages),
            "pages": pages,
        }

    def _extract_text(self, file_path: Path, source_type: str) -> str:
        extractors = {
            "pdf": self._pdf,
            "word": self._word,
            "excel": self._excel,
            "ppt": self._ppt,
        }
        return extractors.get(source_type, self._plain)(file_path)

    def _pdf(self, path: Path) -> str:
        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(str(path))
            chunks = []
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ""
                chunks.append(f"[第{i}页]\n{text}")
            return "\n\n".join(chunks)
        except Exception as e:
            log.warning("PDF extraction failed: %s", e)
            return ""

    def _word(self, path: Path) -> str:
        try:
            import docx
            return "\n".join(p.text for p in docx.Document(str(path)).paragraphs)
        except Exception as e:
            log.warning("Word extraction failed: %s", e)
            return ""

    def _excel(self, path: Path) -> str:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(path), data_only=True)
            rows = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    rows.append("\t".join(str(c) for c in row if c is not None))
            return "\n".join(rows)
        except Exception as e:
            log.warning("Excel extraction failed: %s", e)
            return ""

    def _ppt(self, path: Path) -> str:
        suffix = path.suffix.lower()
        if suffix == ".pptx":
            try:
                from pptx import Presentation

                prs = Presentation(str(path))
                chunks = []
                for i, slide in enumerate(prs.slides, 1):
                    texts = []
                    for shape in slide.shapes:
                        text = getattr(shape, "text", "") or ""
                        text = text.strip()
                        if text:
                            texts.append(text)
                    if texts:
                        chunks.append(f"[第{i}页]\n" + "\n".join(texts))
                return "\n\n".join(chunks)
            except Exception as e:
                log.warning("PPTX extraction failed: %s", e)
                return ""
        log.warning("Legacy/binary PPT (%s); upload OK but text extraction skipped", path.name)
        return ""

    def _plain(self, path: Path) -> str:
        return path.read_text(encoding="utf-8", errors="ignore")
